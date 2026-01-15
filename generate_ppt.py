from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def rgb(r, g, b):
    """Create RGB color for python-pptx 1.0+"""
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

# Color scheme
PRIMARY = rgb(14, 165, 164)
DARK_BG = rgb(15, 23, 42)
CARD_BG = rgb(30, 41, 59)
TEXT_WHITE = rgb(241, 245, 249)
TEXT_MUTED = rgb(148, 163, 184)
PURPLE = rgb(167, 139, 250)

def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_icon(slide):
    icon = slide.shapes.add_textbox(Inches(12.3), Inches(0.3), Inches(0.8), Inches(0.5))
    p = icon.text_frame.paragraphs[0]
    p.text = "+"
    p.font.size = Pt(36)
    p.font.color.rgb = PRIMARY
    p.font.bold = True

def add_pg(slide, n):
    pg = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    p = pg.text_frame.paragraphs[0]
    p.text = str(n)
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

def add_title_text(slide, text, y=0.4, size=36, color=None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(11), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color if color else PRIMARY
    p.font.bold = True

def add_table(slide, data, left, top, widths):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(sum(widths)), Inches(rows * 0.45)).table
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY if ri == 0 else CARD_BG
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(13)
                para.font.color.rgb = TEXT_WHITE
                if ri == 0:
                    para.font.bold = True

# SLIDE 1: Title
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1)
add_icon(s1)
add_pg(s1, 1)

title = s1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1))
p = title.text_frame.paragraphs[0]
p.text = "AI Patient Simulator"
p.font.size = Pt(54)
p.font.color.rgb = PRIMARY
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

sub = s1.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12), Inches(0.6))
p = sub.text_frame.paragraphs[0]
p.text = "An Interactive Medical Training Platform"
p.font.size = Pt(24)
p.font.color.rgb = TEXT_WHITE
p.alignment = PP_ALIGN.CENTER

auth = s1.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(0.5))
p = auth.text_frame.paragraphs[0]
p.text = "Presented by: Amit Kumar"
p.font.size = Pt(20)
p.font.color.rgb = TEXT_MUTED
p.alignment = PP_ALIGN.CENTER

# SLIDE 2: Architecture
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2)
add_icon(s2)
add_pg(s2, 2)
add_title_text(s2, "System Architecture")

def arch_box(slide, txt, sub, top):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(top), Inches(10), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = PRIMARY
    box.line.width = Pt(2)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = sub
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_MUTED
    p2.alignment = PP_ALIGN.CENTER

arch_box(s2, "FRONTEND (Vercel)", "React + Vite | Single Page Application | Modern UI", 1.8)
arr1 = s2.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.2), Inches(2.8), Inches(0.5), Inches(0.4))
arr1.fill.solid()
arr1.fill.fore_color.rgb = PRIMARY
arr1.line.fill.background()

arch_box(s2, "BACKEND (Render)", "Flask + LangChain | Session Management | Auth System", 3.3)
arr2 = s2.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.2), Inches(4.3), Inches(0.5), Inches(0.4))
arr2.fill.solid()
arr2.fill.fore_color.rgb = PRIMARY
arr2.line.fill.background()

arch_box(s2, "LLM (Groq Cloud)", "Llama 3.1 70B (Chat) | Llama 3.1 8B (Generation)", 4.8)

note = s2.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12), Inches(0.5))
p = note.text_frame.paragraphs[0]
p.text = "Deployed separately for scalability | CORS-enabled for cross-origin requests"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p.alignment = PP_ALIGN.CENTER

# SLIDE 3: Technologies
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3)
add_icon(s3)
add_pg(s3, 3)
add_title_text(s3, "Technologies Used")

tech = [
    ["Layer", "Technology", "Purpose"],
    ["Frontend", "React 18 + Vite", "Fast, modern UI framework"],
    ["Styling", "Vanilla CSS", "Glassmorphism dark theme"],
    ["Backend", "Flask (Python)", "REST API server"],
    ["AI Framework", "LangChain", "LLM orchestration"],
    ["LLM Provider", "Groq (Llama 3.1)", "Ultra-fast inference"],
    ["Hosting", "Vercel + Render", "Serverless + containerized"]
]
add_table(s3, tech, 0.5, 1.6, [2, 3, 6])

why = s3.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1.5))
tf = why.text_frame
p = tf.paragraphs[0]
p.text = "Key Benefits:"
p.font.size = Pt(16)
p.font.color.rgb = PRIMARY
p.font.bold = True
for t in ["Groq provides 10x faster inference", "LangChain enables structured LLM interactions", "Separation allows independent scaling"]:
    p = tf.add_paragraph()
    p.text = "  - " + t
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE

# SLIDE 4: Design Choices
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4)
add_icon(s4)
add_pg(s4, 4)
add_title_text(s4, "Key Design Choices")

design = [
    ["Choice", "Implementation"],
    ["Dual-Model Strategy", "8B for generation (fast), 70B for chat (quality)"],
    ["Multi-Key Fallback", "Auto-rotates API keys on 429 rate limit errors"],
    ["Real-Time Symptom Extraction", "LLM extracts symptoms from each response"],
    ["Session State", "In-memory Python dictionaries, no database"],
    ["Privacy-First", "No chat transcripts stored, only summaries"]
]
add_table(s4, design, 0.5, 1.6, [4, 7])

# SLIDE 5: Demo Flow
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5)
add_icon(s5)
add_pg(s5, 5)
add_title_text(s5, "Chatbot Functionality")

flow = [
    ["Step", "Action", "Description"],
    ["1", "Login/Signup", "Doctor authentication with unique username"],
    ["2", "Start Session", "LLM generates randomized patient case"],
    ["3", "Diagnosis", "Doctor asks questions, patient responds"],
    ["4", "Tracking", "Symptoms extracted and displayed in sidebar"],
    ["5", "Analysis", "AI suggests possible conditions"],
    ["6", "End Session", "Doctor enters diagnosis and prescriptions"]
]
add_table(s5, flow, 0.5, 1.6, [1, 2.5, 7.5])

# SLIDE 6: Storage & Features
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s6)
add_icon(s6)
add_pg(s6, 6)
add_title_text(s6, "User Storage & Unique Features")

storage = s6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(3))
tf = storage.text_frame
p = tf.paragraphs[0]
p.text = "In-Memory Python Dictionary"
p.font.size = Pt(18)
p.font.color.rgb = PRIMARY
p.font.bold = True

for t in ["", "No external database required", "Key-value storage structure", "Stores: username, password, history", "", "Trade-offs:", "- Simplifies deployment", "- Data resets on restart", "- Suitable for demo/evaluation"]:
    p = tf.add_paragraph()
    p.text = t
    p.font.size = Pt(13)
    if "Trade" in t:
        p.font.color.rgb = PURPLE
        p.font.bold = True
    else:
        p.font.color.rgb = TEXT_WHITE

feats = [
    ["Feature", "Description"],
    ["True Randomization", "Time-seeded for diverse patients"],
    ["Multi-Key Fallback", "Auto-rotates on rate limit"],
    ["Symptom Tracking", "Real-time extraction"],
    ["AI Analysis", "LLM differential diagnosis"],
    ["Privacy-First", "No chat storage"],
    ["Auth System", "No external database"]
]
add_table(s6, feats, 6, 1.5, [3, 4])

# Save
out = r"c:\Users\amitk\OneDrive\Desktop\ICAPP\AI_Patient_Simulator.pptx"
prs.save(out)
print(f"Saved: {out}")
